from docx import Document
from pptx import Presentation
from typing import List
import os

def export_docx(project_title: str, sections: List[dict], out_path: str):
    """
    Build a .docx file using python-docx.
    sections = [{ "title": str, "content": str }]
    """
    doc = Document()

    doc.add_heading(project_title, level=1)

    for s in sections:
        doc.add_heading(s["title"], level=2)
        doc.add_paragraph(s.get("content", ""))

    doc.save(out_path)
    return out_path


def export_pptx(project_title: str, slides: List[dict], out_path: str):
    """
    Build a .pptx file using python-pptx.
    slides = [{ "title": str, "content": str }]
    """
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = project_title

    # Content slides
    for s in slides:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = s["title"]

        # usually body placeholder is index 1
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.text = s.get("content", "")

    prs.save(out_path)
    return out_path
