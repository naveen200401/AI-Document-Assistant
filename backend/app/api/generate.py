# backend/app/api/generate.py
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
import tempfile, os, uuid

# adjust these imports to your codebase
from app.api.dependencies import get_current_user  # existing auth dependency
# If your LLM client is at app.services.llm_client, import generate_full_document from there
from app.services.llm_client import generate_full_document

router = APIRouter()  # <--- must exist and be exported

class GenerateRequest(BaseModel):
    prompt: str
    format: str = "docx"  # "docx" or "pptx"

def build_docx_from_text(text: str, out_path: str):
    from docx import Document
    doc = Document()
    for block in [b.strip() for b in text.split("\n\n") if b.strip()]:
        doc.add_paragraph(block)
    doc.save(out_path)

def build_pptx_from_text(text: str, out_path: str):
    from pptx import Presentation
    prs = Presentation()
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        prs.slides.add_slide(prs.slide_layouts[5])
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = paragraphs[0][:60]
        if len(paragraphs) > 1:
            body_slide = prs.slides.add_slide(prs.slide_layouts[1])
            tf = body_slide.shapes.placeholders[1].text_frame
            for p in paragraphs[1:]:
                for line in p.split("\n"):
                    txt = line.strip()
                    if txt:
                        tf.add_paragraph().text = txt[:900]
    prs.save(out_path)

@router.post("/generate", response_model=None)
async def generate_document(data: GenerateRequest, current_user=Depends(get_current_user)):
    if data.format not in ("docx", "pptx"):
        raise HTTPException(status_code=400, detail="format must be 'docx' or 'pptx'")

    try:
        text = await generate_full_document(data.prompt)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM generation failed: {e}")

    tmp_dir = tempfile.gettempdir()
    filename = f"generated_{uuid.uuid4().hex}.{data.format}"
    out_path = os.path.join(tmp_dir, filename)

    if data.format == "docx":
        build_docx_from_text(text, out_path)
    else:
        build_pptx_from_text(text, out_path)

    return FileResponse(out_path, filename=filename, media_type=(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if data.format == "docx" else
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ))
