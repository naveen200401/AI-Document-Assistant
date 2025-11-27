"""
backend/app.py

Flask backend for AI Document Assistant.

- Documents with per-user ownership (owner_email)
- Sections (pages) with refine / regenerate / feedback / comments
- Page-by-page generation via Gemini
- Export: DOCX / PDF / PPTX
"""

from __future__ import annotations

import os
import uuid
import sqlite3
from datetime import datetime
from typing import Optional, Dict, Any, List

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS

# -------------------------------------------------
# Flask app + CORS
# -------------------------------------------------
app = Flask(__name__)
CORS(app)

# -------------------------------------------------
# Gemini SDK
# -------------------------------------------------
GENAI_AVAILABLE = False
genai = None
try:
    from google import generativeai as genai  # type: ignore

    GENAI_AVAILABLE = True
except Exception:
    genai = None
    GENAI_AVAILABLE = False

# -------------------------------------------------
# Optional export libs
# -------------------------------------------------
PPTX_AVAILABLE = False
DOCX_AVAILABLE = False
PDF_AVAILABLE = False

try:
    from pptx import Presentation  # type: ignore

    PPTX_AVAILABLE = True
except Exception:
    pass

try:
    from docx import Document  # type: ignore

    DOCX_AVAILABLE = True
except Exception:
    pass

try:
    from reportlab.lib.pagesizes import A4  # type: ignore
    from reportlab.pdfgen import canvas

    PDF_AVAILABLE = True
except Exception:
    pass

# -------------------------------------------------
# Config
# -------------------------------------------------
DB_PATH = os.environ.get("AI_DOCS_DB", "app.db")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
GEMINI_MODEL = os.environ.get("GEMINI_MODEL", "models/gemini-2.5-flash")


# -------------------------------------------------
# DB helpers
# -------------------------------------------------
def db_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    """
    Create tables if missing AND ensure 'owner_email' column exists.
    This is important on Render where an older DB might not have it.
    """
    conn = db_conn()
    cur = conn.cursor()

    # Base tables
    cur.executescript(
        """
        CREATE TABLE IF NOT EXISTS documents (
          id INTEGER PRIMARY KEY AUTOINCREMENT,
          title TEXT,
          owner_email TEXT,
          created_at TEXT DEFAULT (datetime('now'))
        );

        CREATE TABLE IF NOT EXISTS sections (
          id INTEGER PRIMARY KEY AUTOINCREMENT,
          document_id INTEGER NOT NULL,
          position INTEGER DEFAULT 0,
          heading TEXT,
          type TEXT DEFAULT 'text',
          text TEXT,
          last_feedback INTEGER,
          FOREIGN KEY(document_id) REFERENCES documents(id)
        );

        CREATE TABLE IF NOT EXISTS refinements (
          id TEXT PRIMARY KEY,
          section_id INTEGER,
          prompt TEXT,
          revised_text TEXT,
          created_at TEXT DEFAULT (datetime('now')),
          FOREIGN KEY(section_id) REFERENCES sections(id)
        );

        CREATE TABLE IF NOT EXISTS comments (
          id TEXT PRIMARY KEY,
          section_id INTEGER,
          comment TEXT,
          created_at TEXT DEFAULT (datetime('now')),
          FOREIGN KEY(section_id) REFERENCES sections(id)
        );

        CREATE TABLE IF NOT EXISTS feedback (
          id TEXT PRIMARY KEY,
          section_id INTEGER,
          liked INTEGER,
          created_at TEXT DEFAULT (datetime('now')),
          FOREIGN KEY(section_id) REFERENCES sections(id)
        );
        """
    )

    # Safety: if an OLD "documents" table exists without owner_email, add it.
    cur.execute("PRAGMA table_info(documents)")
    cols = [row[1] for row in cur.fetchall()]
    if "owner_email" not in cols:
        cur.execute("ALTER TABLE documents ADD COLUMN owner_email TEXT")

    conn.commit()
    conn.close()


# Make sure DB is ready on Render when gunicorn imports app
with app.app_context():
    init_db()


def serialize_document(doc_id: int) -> Optional[Dict[str, Any]]:
    """Return a document + sections + refinements + comments."""
    conn = db_conn()
    doc = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc:
        conn.close()
        return None

    sections = conn.execute(
        "SELECT * FROM sections WHERE document_id = ? ORDER BY position, id",
        (doc_id,),
    ).fetchall()

    sections_out: List[Dict[str, Any]] = []
    for s in sections:
        sdict = dict(s)
        refinements = conn.execute(
            "SELECT * FROM refinements WHERE section_id = ? ORDER BY created_at DESC",
            (s["id"],),
        ).fetchall()
        comments = conn.execute(
            "SELECT * FROM comments WHERE section_id = ? ORDER BY created_at",
            (s["id"],),
        ).fetchall()
        sdict["refinements"] = [dict(r) for r in refinements]
        sdict["comments"] = [dict(c) for c in comments]
        sections_out.append(sdict)

    out = dict(doc)
    out["sections"] = sections_out
    conn.close()
    return out


# -------------------------------------------------
# Gemini helpers
# -------------------------------------------------
def _extract_gemini_text(resp: Any) -> Optional[str]:
    """Safely extract text from a google-generativeai response."""
    # 1) Quick accessor
    try:
        t = getattr(resp, "text", None)
        if isinstance(t, str) and t.strip():
            return t.strip()
    except Exception as e:
        print("resp.text failed:", e)

    # 2) candidates[0].content.parts[*].text
    try:
        candidates = getattr(resp, "candidates", None) or []
        if candidates:
            first = candidates[0]
            finish_reason = getattr(first, "finish_reason", None)
            print("Gemini finish_reason:", finish_reason)

            content = getattr(first, "content", None)
            parts = getattr(content, "parts", None) if content is not None else None
            texts: List[str] = []
            if parts:
                for p in parts:
                    pt = getattr(p, "text", None)
                    if isinstance(pt, str) and pt.strip():
                        texts.append(pt.strip())
            if texts:
                return "\n".join(texts)
    except Exception as e:
        print("candidate.parts extraction failed:", e)

    # 3) fallback to string
    try:
        s = str(resp)
        if s.strip():
            return s.strip()
    except Exception:
        pass

    return None


def call_gemini_text(prompt: str, model_name: Optional[str] = None) -> str:
    if not GENAI_AVAILABLE:
        raise RuntimeError("google-generativeai SDK not installed.")
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY environment variable not set.")

    genai.configure(api_key=GEMINI_API_KEY)
    model_id = model_name or GEMINI_MODEL
    model = genai.GenerativeModel(model_id)
    resp = model.generate_content(prompt)

    text = _extract_gemini_text(resp)
    if not text or not text.strip():
        finish_reason = None
        try:
            candidates = getattr(resp, "candidates", None) or []
            if candidates:
                finish_reason = getattr(candidates[0], "finish_reason", None)
        except Exception:
            pass
        raise RuntimeError(
            f"Gemini generation failed: no usable text (finish_reason={finish_reason})"
        )

    return text.strip()


# -------------------------------------------------
# Local fallbacks
# -------------------------------------------------
def fallback_page_text(page_index: int, total_pages: int, user_prompt: str) -> str:
    idx = page_index + 1
    return (
        f"This is placeholder content for page {idx} of {total_pages} based on the user prompt: "
        f"'{user_prompt}'.\n\n"
        "Gemini could not be reached right now, so this text was generated locally. "
        "You can refine or regenerate this page later once the model is available."
    )


def fallback_refinement_text(base_text: str, prompt: str) -> str:
    p = (prompt or "").lower()
    if not base_text:
        return "[no content]"

    if "shorten" in p:
        import re

        m = re.search(r"(\d+)", p)
        if m:
            n = int(m.group(1))
            words = base_text.split()
            return " ".join(words[:n]) + ("…" if len(words) > n else "")
        words = base_text.split()
        n = max(1, len(words) // 2)
        return " ".join(words[:n]) + ("…" if len(words) > n else "")

    if "bullet" in p or "bullets" in p:
        import re

        sentences = re.split(r"(?<=[.!?])\s+", base_text.strip())
        bullets = "\n".join("- " + s.strip() for s in sentences if s.strip())
        return bullets

    if "formal" in p:
        return "In a more formal tone: " + base_text

    return base_text + f"\n\n[Refined locally with prompt: {prompt}]"


# -------------------------------------------------
# Basic routes
# -------------------------------------------------
@app.route("/")
def root():
    return jsonify(
        {
            "ok": True,
            "message": "AI Docs backend (Flask + Gemini)",
            "genai_available": GENAI_AVAILABLE,
            "pptx_available": PPTX_AVAILABLE,
            "docx_available": DOCX_AVAILABLE,
            "pdf_available": PDF_AVAILABLE,
        }
    )


@app.route("/api/health")
def health():
    return jsonify({"status": "ok"})


# -------------------------------------------------
# Documents list / create (per-user)
# -------------------------------------------------
@app.route("/api/documents", methods=["GET", "POST"])
def documents():
    conn = db_conn()
    cur = conn.cursor()

    if request.method == "POST":
        data = request.get_json(force=True) or {}
        title = data.get("title") or "Untitled"
        owner_email = (data.get("owner_email") or "").strip()

        cur.execute(
            "INSERT INTO documents (title, owner_email) VALUES (?, ?)",
            (title, owner_email),
        )
        doc_id = cur.lastrowid
        conn.commit()

        doc = cur.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
        conn.close()
        return jsonify({**dict(doc), "sections": []}), 201

    # GET – list only this user's documents
    owner_email = (request.args.get("owner_email") or "").strip()

    if owner_email:
        rows = cur.execute(
            """
            SELECT * FROM documents
            WHERE owner_email = ?
            ORDER BY datetime(created_at) DESC, id DESC
            """,
            (owner_email,),
        ).fetchall()
    else:
        # if for some reason email missing, just return empty list
        rows = []

    conn.close()
    return jsonify([dict(r) for r in rows]), 200


@app.route("/api/document/<int:doc_id>", methods=["GET"])
def get_document(doc_id: int):
    owner_email = (request.args.get("owner_email") or "").strip()

    conn = db_conn()
    doc = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc:
        conn.close()
        return jsonify({"error": "document not found"}), 404

    # Simple access control: document must belong to this email
    if owner_email and (doc["owner_email"] or "") != owner_email:
        conn.close()
        return jsonify({"error": "document not found"}), 404

    sections = conn.execute(
        "SELECT * FROM sections WHERE document_id = ? ORDER BY position, id",
        (doc_id,),
    ).fetchall()

    sections_out = []
    for s in sections:
        sdict = dict(s)
        refinements = conn.execute(
            "SELECT * FROM refinements WHERE section_id = ? ORDER BY created_at DESC",
            (s["id"],),
        ).fetchall()
        comments = conn.execute(
            "SELECT * FROM comments WHERE section_id = ? ORDER BY created_at",
            (s["id"],),
        ).fetchall()
        sdict["refinements"] = [dict(r) for r in refinements]
        sdict["comments"] = [dict(c) for c in comments]
        sections_out.append(sdict)

    out = dict(doc)
    out["sections"] = sections_out
    conn.close()
    return jsonify(out), 200


# -------------------------------------------------
# Generate pages
# -------------------------------------------------
@app.route("/api/documents/<int:doc_id>/generate", methods=["POST"])
def generate_document(doc_id: int):
    data = request.get_json(force=True) or {}
    user_prompt = (data.get("prompt") or "").strip()
    theme = (data.get("theme") or "").strip()
    pages = int(data.get("pages") or 1)
    pages = max(1, min(pages, 30))

    if not user_prompt:
        return jsonify({"error": "prompt is required"}), 400

    conn = db_conn()
    doc = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc:
        conn.close()
        return jsonify({"error": "document not found"}), 404

    # Clear existing sections + child tables
    sec_rows = conn.execute(
        "SELECT id FROM sections WHERE document_id = ?", (doc_id,)
    ).fetchall()
    sec_ids = [r["id"] for r in sec_rows]
    if sec_ids:
        q = ",".join("?" for _ in sec_ids)
        conn.execute(f"DELETE FROM refinements WHERE section_id IN ({q})", sec_ids)
        conn.execute(f"DELETE FROM comments    WHERE section_id IN ({q})", sec_ids)
        conn.execute(f"DELETE FROM feedback    WHERE section_id IN ({q})", sec_ids)
        conn.execute("DELETE FROM sections WHERE document_id = ?", (doc_id,))
        conn.commit()

    # Generate each page
    for i in range(pages):
        heading = f"Page {i + 1}"
        page_prompt = (
            "You are a document author. Generate the content for page "
            f"{i+1} of {pages} for the user prompt.\n\n"
            f"USER_PROMPT:\n{user_prompt}\n\n"
            "PAGE_INSTRUCTIONS: Generate one page of text suitable for a slide "
            "or document page. Keep it concise (about 120–200 words) unless "
            "the user requested otherwise."
        )
        if theme:
            page_prompt += f"\nTheme or tone: {theme}."

        try:
            page_text = call_gemini_text(page_prompt)
        except Exception as e:
            print(f"[WARN] Gemini failed for page {i+1}: {e}")
            page_text = fallback_page_text(i, pages, user_prompt)

        cur = conn.cursor()
        cur.execute(
            """
            INSERT INTO sections (document_id, position, heading, type, text)
            VALUES (?, ?, ?, 'text', ?)
            """,
            (doc_id, i, heading, page_text),
        )
        section_id = cur.lastrowid

        rid = str(uuid.uuid4())
        cur.execute(
            """
            INSERT INTO refinements (id, section_id, prompt, revised_text, created_at)
            VALUES (?, ?, ?, ?, ?)
            """,
            (
                rid,
                section_id,
                page_prompt,
                page_text,
                datetime.utcnow().isoformat(),
            ),
        )
        conn.commit()

    conn.close()
    doc_out = serialize_document(doc_id)
    return jsonify(doc_out), 200


# -------------------------------------------------
# Section refine / regenerate / feedback / comment
# -------------------------------------------------
@app.route("/api/sections/<int:section_id>/refine", methods=["POST"])
def refine_section(section_id: int):
    data = request.get_json(force=True) or {}
    prompt = (data.get("prompt") or "").strip()
    current_text = data.get("current_text")

    if not prompt:
        return jsonify({"error": "prompt required"}), 400

    conn = db_conn()
    sec = conn.execute("SELECT * FROM sections WHERE id = ?", (section_id,)).fetchone()
    if not sec:
        conn.close()
        return jsonify({"error": "section not found"}), 404

    base_text = current_text if current_text is not None else (sec["text"] or "")
    full_prompt = (
        "You are an expert document editor. Apply the user's refinement instruction "
        "ONLY to the provided section text.\n\n"
        f"SECTION_TEXT:\n{base_text}\n\n"
        f"REFINEMENT_INSTRUCTION:\n{prompt}\n\n"
        "Return the revised section text only (no extra commentary)."
    )

    try:
        revised_text = call_gemini_text(full_prompt)
    except Exception as e:
        print("Gemini refine failed, using fallback:", e)
        revised_text = fallback_refinement_text(base_text, prompt)

    rid = str(uuid.uuid4())
    created_at = datetime.utcnow().isoformat()
    conn.execute(
        """
        INSERT INTO refinements (id, section_id, prompt, revised_text, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (rid, section_id, prompt, revised_text, created_at),
    )
    conn.execute("UPDATE sections SET text = ? WHERE id = ?", (revised_text, section_id))
    conn.commit()

    refinements = conn.execute(
        "SELECT * FROM refinements WHERE section_id = ? ORDER BY created_at DESC",
        (section_id,),
    ).fetchall()
    conn.close()
    return jsonify(
        {
            "revised_text": revised_text,
            "refinements": [dict(r) for r in refinements],
        }
    )


@app.route("/api/sections/<int:section_id>/regenerate", methods=["POST"])
def regenerate_section(section_id: int):
    conn = db_conn()
    sec = conn.execute("SELECT * FROM sections WHERE id = ?", (section_id,)).fetchone()
    if not sec:
        conn.close()
        return jsonify({"error": "section not found"}), 404

    doc = conn.execute(
        "SELECT * FROM documents WHERE id = ?", (sec["document_id"],)
    ).fetchone()
    doc_title = doc["title"] if doc else "this document"

    regen_prompt = (
        "You are a document author. Regenerate this section as a fresh page while "
        "keeping the topic and intent similar.\n\n"
        f"DOCUMENT_TITLE: {doc_title}\n\n"
        f"CURRENT_SECTION_HEADING: {sec['heading']}\n\n"
        f"CURRENT_SECTION_TEXT:\n{sec['text']}\n\n"
        "Return a single improved page of text (about 120–200 words) and no extra notes."
    )

    try:
        new_text = call_gemini_text(regen_prompt)
    except Exception as e:
        print("Gemini regenerate failed, using fallback:", e)
        new_text = fallback_refinement_text(sec["text"] or "", "regenerate")

    rid = str(uuid.uuid4())
    created_at = datetime.utcnow().isoformat()
    conn.execute(
        """
        INSERT INTO refinements (id, section_id, prompt, revised_text, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (rid, section_id, "[Regenerate page]", new_text, created_at),
    )
    conn.execute("UPDATE sections SET text = ? WHERE id = ?", (new_text, section_id))
    conn.commit()

    refinements = conn.execute(
        "SELECT * FROM refinements WHERE section_id = ? ORDER BY created_at DESC",
        (section_id,),
    ).fetchall()
    conn.close()
    return jsonify(
        {
            "id": section_id,
            "text": new_text,
            "refinements": [dict(r) for r in refinements],
        }
    )


@app.route("/api/sections/<int:section_id>/feedback", methods=["POST"])
def section_feedback(section_id: int):
    data = request.get_json(force=True) or {}
    liked = data.get("liked")
    if liked is None:
        return jsonify({"error": "liked must be provided"}), 400

    conn = db_conn()
    sec = conn.execute("SELECT * FROM sections WHERE id = ?", (section_id,)).fetchone()
    if not sec:
        conn.close()
        return jsonify({"error": "section not found"}), 404

    fid = str(uuid.uuid4())
    created_at = datetime.utcnow().isoformat()
    conn.execute(
        "INSERT INTO feedback (id, section_id, liked, created_at) VALUES (?, ?, ?, ?)",
        (fid, section_id, int(bool(liked)), created_at),
    )
    conn.execute(
        "UPDATE sections SET last_feedback = ? WHERE id = ?",
        (int(bool(liked)), section_id),
    )
    conn.commit()
    conn.close()
    return jsonify({"ok": True})


@app.route("/api/sections/<int:section_id>/comment", methods=["POST"])
def section_comment(section_id: int):
    data = request.get_json(force=True) or {}
    comment_text = (data.get("comment") or "").strip()
    if not comment_text:
        return jsonify({"error": "comment required"}), 400

    conn = db_conn()
    sec = conn.execute("SELECT * FROM sections WHERE id = ?", (section_id,)).fetchone()
    if not sec:
        conn.close()
        return jsonify({"error": "section not found"}), 404

    cid = str(uuid.uuid4())
    created_at = datetime.utcnow().isoformat()
    conn.execute(
        """
        INSERT INTO comments (id, section_id, comment, created_at)
        VALUES (?, ?, ?, ?)
        """,
        (cid, section_id, comment_text, created_at),
    )
    conn.commit()
    new_comment = {
        "id": cid,
        "section_id": section_id,
        "comment": comment_text,
        "created_at": created_at,
    }
    conn.close()
    return jsonify({"comment": new_comment})


# -------------------------------------------------
# Export DOCX / PPTX / PDF
# -------------------------------------------------
def _safe_filename(title: str, ext: str) -> str:
    base = "".join(c for c in title if c.isalnum() or c in (" ", "_", "-")).strip()
    if not base:
        base = "document"
    return f"{base}.{ext}"


@app.route("/api/documents/<int:doc_id>/export", methods=["GET"])
def export_document(doc_id: int):
    fmt = (request.args.get("format") or "docx").lower()
    doc = serialize_document(doc_id)
    if not doc:
        return jsonify({"error": "document not found"}), 404

    title = doc.get("title") or f"Document_{doc_id}"
    sections = doc.get("sections") or []

    # DOCX
    if fmt == "docx":
        if not DOCX_AVAILABLE:
            return jsonify(
                {"error": "DOCX export not available (python-docx missing)"}
            ), 500

        from tempfile import NamedTemporaryFile

        d = Document()
        d.add_heading(title, level=1)
        for s in sections:
            d.add_heading(s.get("heading") or "", level=2)
            d.add_paragraph(s.get("text") or "")
        tmp = NamedTemporaryFile(delete=False, suffix=".docx")
        d.save(tmp.name)
        tmp.flush()
        filename = _safe_filename(title, "docx")
        return send_file(
            tmp.name,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    # PPTX
    if fmt == "pptx":
        if not PPTX_AVAILABLE:
            return jsonify(
                {"error": "PPTX export not available (python-pptx missing)"}
            ), 500

        from tempfile import NamedTemporaryFile
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import re

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Generated with AI Document Assistant"
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(0, 70, 140)

        # Content slides
        content_layout = prs.slide_layouts[1]
        for idx, s in enumerate(sections):
            heading = s.get("heading") or f"Page {idx + 1}"
            text = (s.get("text") or "").strip()
            if not text:
                continue

            slide = prs.slides.add_slide(content_layout)

            slide_title = slide.shapes.title
            slide_title.text = heading
            for p in slide_title.text_frame.paragraphs:
                p.font.size = Pt(30)
                p.font.color.rgb = RGBColor(0, 70, 140)

            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            chunks = re.split(r"(?<=[.!?])\s+", text)
            first = True
            for chunk in chunks:
                chunk = chunk.strip()
                if not chunk:
                    continue

                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = chunk
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(40, 40, 40)

        tmp = NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(tmp.name)
        tmp.flush()
        filename = _safe_filename(title, "pptx")
        return send_file(
            tmp.name,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    # PDF
    if fmt == "pdf":
        if not PDF_AVAILABLE:
            return jsonify(
                {"error": "PDF export not available (reportlab missing)"}
            ), 500

        from tempfile import NamedTemporaryFile

        tmp = NamedTemporaryFile(delete=False, suffix=".pdf")
        c = canvas.Canvas(tmp.name, pagesize=A4)
        width, height = A4
        margin = 50
        y = height - margin

        c.setFont("Helvetica-Bold", 16)
        c.drawString(margin, y, title)
        y -= 30

        c.setFont("Helvetica", 11)
        for s in sections:
            if y < margin + 50:
                c.showPage()
                y = height - margin
                c.setFont("Helvetica", 11)

            c.setFont("Helvetica-Bold", 12)
            c.drawString(margin, y, s.get("heading") or "")
            y -= 18
            c.setFont("Helvetica", 11)
            for line in (s.get("text") or "").splitlines():
                if y < margin + 40:
                    c.showPage()
                    y = height - margin
                    c.setFont("Helvetica", 11)
                c.drawString(margin, y, line[:1000])
                y -= 14
            y -= 10

        c.showPage()
        c.save()
        tmp.flush()
        filename = _safe_filename(title, "pdf")
        return send_file(
            tmp.name,
            as_attachment=True,
            download_name=filename,
            mimetype="application/pdf",
        )

    return jsonify({"error": f"unknown format '{fmt}'"}), 400


# -------------------------------------------------
# Debug Gemini
# -------------------------------------------------
@app.route("/api/debug_gemini", methods=["POST"])
def debug_gemini():
    data = request.get_json(force=True) or {}
    prompt = (data.get("prompt") or "").strip()
    model = (data.get("model") or GEMINI_MODEL).strip()

    if not prompt:
        return jsonify({"error": "prompt required"}), 400

    if not GENAI_AVAILABLE or not GEMINI_API_KEY:
        return jsonify({"error": "Gemini SDK or API key not available"}), 500

    try:
        text = call_gemini_text(prompt, model_name=model)
        return jsonify({"prompt": prompt, "model": model, "text": text})
    except Exception as e:
        return jsonify({"error": "Gemini generation failed", "details": {"message": str(e)}}), 500


# -------------------------------------------------
# Main entry (local dev)
# -------------------------------------------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", "5001"))
    print(
        f"Starting AI Docs backend (Gemini) on port {port} "
        f"(GENAI_AVAILABLE={GENAI_AVAILABLE}, "
        f"PPTX_AVAILABLE={PPTX_AVAILABLE}, DOCX_AVAILABLE={DOCX_AVAILABLE}, "
        f"PDF_AVAILABLE={PDF_AVAILABLE})"
    )
    app.run(host="0.0.0.0", port=port, debug=True)
